#!/usr/bin/env python3
"""Build the WSQ Agile Project Management for Business slide deck (all-white house style).

Design grammar and helpers are the proven v2 reference set (cover, section, content,
two_col, cards3, tile_grid, flow_h, process_map, decision_map, compare_table,
trainer_slide, activity_overview, test_slide, brk) EXTENDED here with:

  * img_points / img_full  — place the generated diagram + chart assets
  * tool_slide             — the browser-mock ed-tool slide for each activity's tool
  * activity_slide         — ONE slide per activity (no step-by-step in the deck)

HOUSE RULES APPLIED HERE (per the course owner's brief):
  * Step-by-step procedures live ONLY in the Learner Guide — this deck shows the
    activity briefing, the tool and the expected outcome, never the numbered steps.
  * NO practice-exam slide.
  * Every asset in courseware/assets/ must land on a slide (asserted at the end).

Content is driven entirely by course_data.py + data_domain{1,2,3}.py.
"""
import os, sys, copy, math, re, json
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn
from lxml import etree
from PIL import Image

HERE = os.path.dirname(os.path.abspath(__file__))
sys.path.insert(0, HERE)
import course_data as C
from data_domain1 import DOMAIN1
from data_domain2 import DOMAIN2
from data_domain3 import DOMAIN3
ACTIVITIES = DOMAIN1 + DOMAIN2 + DOMAIN3


def _find_repo(start):
    env = os.environ.get("COURSE_REPO")
    if env and os.path.isdir(env):
        return env
    d = start
    for _ in range(8):
        d = os.path.dirname(d)
        if os.path.isdir(os.path.join(d, "courseware")) and os.path.isdir(os.path.join(d, "activities")):
            return d
    return os.path.dirname(os.path.dirname(HERE))


REPO = _find_repo(HERE)
SKILL_ASSETS = os.path.join(os.path.dirname(HERE), "assets")
CHARTS = os.path.join(REPO, "courseware", "assets")
USED_ASSETS = set()

# the real tool screenshot used on each activity's tool slide
TOOL_SHOT = {
    "Design Thinking": "tool-designthinking.png",
    "Fishbone": "tool-fishbone.png",
    "Scrum Board": "tool-scrum.png",
    "RACI Matrix": "tool-raci.png",
    "5 Whys": "tool-5whys.png",
    "Pareto Chart": "tool-paretochart.png",
}

# ---------------- palette ----------------
BLUE = RGBColor(0x1F, 0x6F, 0xEB); TEAL = RGBColor(0x10, 0xB9, 0x81)
AMBER = RGBColor(0xF5, 0x9E, 0x0B); INK = RGBColor(0x16, 0x1B, 0x26)
GREY = RGBColor(0x5B, 0x63, 0x72); LIGHT = RGBColor(0xF5, 0xF8, 0xFC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF); LINE = RGBColor(0xE2, 0xE8, 0xF0)
VIOLET = RGBColor(0x7C, 0x3A, 0xED); RED = RGBColor(0xDC, 0x26, 0x26)
GREEN = RGBColor(0x12, 0x7A, 0x3E); NAVY = RGBColor(0x0B, 0x12, 0x20)
PALETTE = [BLUE, TEAL, VIOLET, AMBER]

prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]
PAGE = {"n": 1}
SLIDE_MAP = {}


def mark(key):
    """Record the CURRENT slide number for a section/activity so the LP can cite it."""
    SLIDE_MAP[key] = PAGE["n"] + 1


# ---------------- primitives ----------------
def slide(): return prs.slides.add_slide(BLANK)


def rect(s, x, y, w, h, color, line=None):
    sp = s.shapes.add_shape(1, x, y, w, h); sp.fill.solid(); sp.fill.fore_color.rgb = color
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb = line; sp.line.width = Pt(1)
    sp.shadow.inherit = False; return sp


def oval(s, x, y, w, h, color):
    sp = s.shapes.add_shape(MSO_SHAPE.OVAL, x, y, w, h); sp.fill.solid()
    sp.fill.fore_color.rgb = color; sp.line.fill.background(); sp.shadow.inherit = False; return sp


def roundrect(s, x, y, w, h, color, line=None, adj=0.10):
    sp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb = line; sp.line.width = Pt(1.4)
    try: sp.adjustments[0] = adj
    except Exception: pass
    sp.shadow.inherit = False; return sp


def txt(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space=4):
    tb = s.shapes.add_textbox(x, y, w, h); tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, line in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(space)
        for t, sz, col, bold in line:
            r = p.add_run(); r.text = t; r.font.size = Pt(sz); r.font.bold = bold
            r.font.color.rgb = col; r.font.name = "Arial"
    return tb


def bullets(s, x, y, w, h, items, size=18, color=INK, gap=10, mcolor=BLUE):
    tb = s.shapes.add_textbox(x, y, w, h); tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        r = p.add_run(); r.text = "•  "; r.font.size = Pt(size); r.font.bold = True
        r.font.color.rgb = mcolor; r.font.name = "Arial"
        r2 = p.add_run(); r2.text = it; r2.font.size = Pt(size)
        r2.font.color.rgb = color; r2.font.name = "Arial"
    return tb


def connector(s, x1, y1, x2, y2, color, width=Pt(2.0), arrow=True):
    cn = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    cn.line.color.rgb = color; cn.line.width = width
    if arrow:
        ln = cn.line._get_or_add_ln()
        tail = etree.SubElement(ln, qn('a:tailEnd'))
        tail.set('type', 'triangle'); tail.set('w', 'med'); tail.set('len', 'med')
    return cn


# ---------------- motion ----------------
_MC_NS = "http://schemas.openxmlformats.org/markup-compatibility/2006"
_P14_NS = "http://schemas.microsoft.com/office/powerpoint/2010/main"
_P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"


def _transition(s, kind="fade", speed="med"):
    """Restrained slide transition. Written as a raw mc:AlternateContent block so it
    round-trips through PowerPoint; python-pptx does not register the 'mc' prefix,
    so the namespaces are spelled out in full here."""
    xml = s._element
    for old in xml.findall("{%s}AlternateContent" % _MC_NS):
        xml.remove(old)
    ac = etree.SubElement(xml, "{%s}AlternateContent" % _MC_NS,
                          nsmap={"mc": _MC_NS})
    ch = etree.SubElement(ac, "{%s}Choice" % _MC_NS, nsmap={"p14": _P14_NS})
    ch.set("Requires", "p14")
    tr = etree.SubElement(ch, "{%s}transition" % _P_NS)
    tr.set("spd", speed)
    tr.set("{%s}dur" % _P14_NS, "700" if kind == "fade" else "900")
    if kind == "fade":
        etree.SubElement(tr, "{%s}fade" % _P_NS)
    elif kind == "push":
        el = etree.SubElement(tr, "{%s}push" % _P_NS); el.set("dir", "u")
    else:
        el = etree.SubElement(tr, "{%s}wipe" % _P_NS); el.set("dir", "r")


# ---------------- text fitting ----------------
def _ellipsis(text, limit):
    t = " ".join(str(text).split())
    if len(t) <= limit: return t
    cut = t[:limit]; sp = cut.rfind(" ")
    if sp > limit * 0.55: cut = cut[:sp]
    return cut.rstrip(" ,.;:-—(") + "…"


def _fit_title(title, size=29):
    n = len(title)
    if n <= 52: return size
    if n <= 66: return 25
    if n <= 82: return 22
    return 20


def _fit_pt(text, width_in, height_in, base=12.0, floor=8.0):
    """Return the largest point size (<= base, >= floor) at which `text` still fits
    inside a width_in x height_in box.

    Text on a slide must never be truncated — an ellipsis silently discards
    instructional content. We shrink the type instead, using the empirical rule that
    Arial averages ~0.50 em per character and needs ~1.22 em of leading per line.
    """
    t = " ".join(str(text).split())
    if not t:
        return base
    size = base
    while size > floor:
        chars_per_line = max(1, int((width_in * 72.0) / (size * 0.50)))
        lines = -(-len(t) // chars_per_line)          # ceil
        if lines * (size * 1.22) <= height_in * 72.0:
            return round(size, 1)
        size -= 0.5
    return floor


def footer(s):
    PAGE["n"] += 1
    txt(s, Inches(0.4), Inches(7.05), Inches(7.5), Inches(0.35),
        [[(f"{C.SHORT_TITLE}  ·  {C.COURSE_CODE}", 9, GREY, False)]])
    txt(s, Inches(5.0), Inches(7.05), Inches(3.3), Inches(0.35),
        [[("© 2026 Tertiary Infotech Academy Pte Ltd", 9, GREY, False)]], align=PP_ALIGN.CENTER)
    txt(s, Inches(12.4), Inches(7.05), Inches(0.6), Inches(0.35),
        [[(str(PAGE["n"]), 9, GREY, False)]], align=PP_ALIGN.RIGHT)


def head(s, title, kicker=None, kcolor=BLUE):
    rect(s, 0, 0, SW, SH, WHITE); rect(s, 0, 0, Inches(0.28), Inches(1.55), kcolor)
    if kicker:
        txt(s, Inches(0.85), Inches(0.5), Inches(11.6), Inches(0.4), [[(kicker, 14, kcolor, True)]])
    txt(s, Inches(0.85), Inches(0.88), Inches(11.9), Inches(0.78),
        [[(title, _fit_title(title), INK, True)]], anchor=MSO_ANCHOR.MIDDLE)
    rect(s, Inches(0.85), Inches(1.7), Inches(11.63), Inches(0.02), LINE)
    return s


def _asset(name):
    p = os.path.join(CHARTS, name)
    if os.path.exists(p):
        USED_ASSETS.add(name); return p
    p2 = os.path.join(SKILL_ASSETS, name)
    return p2 if os.path.exists(p2) else None


def _fit(path, maxw_in, maxh_in):
    """Aspect-fit an image into a box, returning (w,h) in EMU."""
    with Image.open(path) as im:
        iw, ih = im.size
    sc = min(maxw_in / iw, maxh_in / ih)
    return Inches(iw * sc), Inches(ih * sc)


# ============================================================ slide templates
def cover():
    s = slide(); rect(s, 0, 0, SW, SH, WHITE)
    rect(s, 0, 0, SW, Inches(0.22), BLUE); rect(s, 0, Inches(7.28), SW, Inches(0.22), TEAL)
    org = _asset("tertiary-infotech-logo.png")
    if org: s.shapes.add_picture(org, Inches(0.85), Inches(0.7), height=Inches(1.05))
    # course badge (top-right)
    rect(s, Inches(10.45), Inches(0.66), Inches(2.15), Inches(1.16), BLUE)
    txt(s, Inches(10.45), Inches(0.80), Inches(2.15), Inches(0.5), [[("AGILE", 22, WHITE, True)]],
        align=PP_ALIGN.CENTER)
    txt(s, Inches(10.45), Inches(1.34), Inches(2.15), Inches(0.4), [[("PROJECT MGMT", 10, WHITE, True)]],
        align=PP_ALIGN.CENTER)
    txt(s, Inches(0.9), Inches(2.3), Inches(12), Inches(0.6), [[("TRAINER SLIDES  ·  WSQ", 16, BLUE, True)]])
    txt(s, Inches(0.9), Inches(2.85), Inches(12.0), Inches(1.9), [[(C.TITLE, 40, INK, True)]])
    rect(s, Inches(0.92), Inches(4.75), Inches(2.4), Inches(0.06), TEAL)
    txt(s, Inches(0.9), Inches(5.05), Inches(12), Inches(1.4),
        [[(f"WSQ Course Code: {C.COURSE_CODE}", 16, GREY, False)],
         [(f"Conducted by {C.ORG}  ·  {C.UEN}", 14, GREY, False)],
         [(f"{C.DAYS} days  ·  {C.DURATION_HRS} training hours  ·  TSC: {C.TSC_TITLE} ({C.TSC_CODE})", 13, GREY, False)]],
        space=6)
    txt(s, Inches(0.9), Inches(6.55), Inches(12), Inches(0.4),
        [[(f"Version {C.VERSION}  ·  {C.VERSION_DATE}", 12, GREY, False)]])
    txt(s, Inches(0.9), Inches(6.88), Inches(12), Inches(0.34),
        [[("© 2026 Tertiary Infotech Academy Pte Ltd. All rights reserved.  ·  www.tertiarycourses.com.sg",
           10, GREY, False)]])


def section(kicker, title, n, sub=""):
    s = slide(); rect(s, 0, 0, SW, SH, WHITE)
    rect(s, Inches(0.72), Inches(1.90), Inches(11.85), Inches(3.60), LIGHT)
    rect(s, Inches(0.72), Inches(1.90), Inches(0.14), Inches(3.60), BLUE)
    txt(s, Inches(1.40), Inches(2.35), Inches(11), Inches(0.5), [[(kicker, 15, BLUE, True)]])
    txt(s, Inches(1.40), Inches(2.85), Inches(11.0), Inches(1.05), [[(title, 34, INK, True)]])
    if sub:
        txt(s, Inches(1.40), Inches(3.95), Inches(10.8), Inches(1.2), [[(sub, 15, GREY, False)]])
    if n:
        txt(s, Inches(10.0), Inches(0.55), Inches(2.6), Inches(1.5),
            [[(n, 72, RGBColor(0xE2, 0xE8, 0xF0), True)]], align=PP_ALIGN.RIGHT)
    footer(s); return s


def content(title, items, kicker=None, size=19):
    s = head(slide(), title, kicker)
    bullets(s, Inches(0.85), Inches(1.95), Inches(11.6), Inches(4.9), items, size=size)
    footer(s); return s


def two_col(title, left, right, kicker=None, lhead="", rhead="", note=None):
    s = head(slide(), title, kicker)
    bh = Inches(4.7) if not note else Inches(3.95)
    rect(s, Inches(0.85), Inches(1.95), Inches(5.7), bh, LIGHT)
    rect(s, Inches(6.78), Inches(1.95), Inches(5.7), bh, LIGHT)
    rect(s, Inches(0.85), Inches(1.95), Inches(5.7), Inches(0.42), RED)
    rect(s, Inches(6.78), Inches(1.95), Inches(5.7), Inches(0.42), TEAL)
    if lhead:
        txt(s, Inches(0.85), Inches(1.95), Inches(5.7), Inches(0.42), [[(lhead, 14, WHITE, True)]],
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    if rhead:
        txt(s, Inches(6.78), Inches(1.95), Inches(5.7), Inches(0.42), [[(rhead, 14, WHITE, True)]],
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    bullets(s, Inches(1.1), Inches(2.6), Inches(5.2), bh - Inches(0.8), left, size=14, mcolor=RED, gap=8)
    bullets(s, Inches(7.03), Inches(2.6), Inches(5.2), bh - Inches(0.8), right, size=14, mcolor=TEAL, gap=8)
    if note:
        rect(s, Inches(0.85), Inches(6.08), Inches(11.63), Inches(0.72), LIGHT)
        rect(s, Inches(0.85), Inches(6.08), Inches(0.1), Inches(0.72), VIOLET)
        txt(s, Inches(1.15), Inches(6.08), Inches(11.2), Inches(0.72), [[(note, 12.5, INK, True)]],
            anchor=MSO_ANCHOR.MIDDLE)
    footer(s); return s


def cards3(title, cards, kicker):
    s = head(slide(), title, kicker); xs = [Inches(0.85), Inches(4.81), Inches(8.77)]
    for i, c in enumerate(cards[:3]):
        x = xs[i]; col = c[0]
        rect(s, x, Inches(1.95), Inches(3.71), Inches(4.7), LIGHT)
        rect(s, x, Inches(1.95), Inches(3.71), Inches(0.12), col)
        txt(s, x + Inches(0.25), Inches(2.2), Inches(3.25), Inches(0.6), [[(c[1], 18, col, True)]])
        bullets(s, x + Inches(0.25), Inches(2.95), Inches(3.25), Inches(3.4), c[2], size=13,
                mcolor=col, gap=8)
    footer(s); return s


def big_statement(line1, line2, kicker, color=BLUE):
    s = slide(); rect(s, 0, 0, SW, SH, WHITE); rect(s, 0, 0, Inches(0.28), SH, color)
    txt(s, Inches(1.1), Inches(2.2), Inches(11), Inches(0.5), [[(kicker, 16, color, True)]])
    txt(s, Inches(1.1), Inches(2.8), Inches(11.3), Inches(2.4), [[(line1, 36, INK, True)]])
    if line2:
        txt(s, Inches(1.12), Inches(4.95), Inches(11), Inches(1.3), [[(line2, 19, GREY, False)]])
    footer(s); return s


def tile_grid(title, items, kicker=None, cols=2, size=15, icons=None, accent=BLUE):
    s = head(slide(), title, kicker, kcolor=accent)
    n = len(items); rows = math.ceil(n / cols)
    X0 = Inches(0.85); Y0 = Inches(1.95); TOTW = Inches(11.63); AREAH = Inches(4.78)
    gx = Inches(0.3); gy = Inches(0.26)
    cw = int((TOTW - gx * (cols - 1)) / cols); ch = int((AREAH - gy * (rows - 1)) / rows)
    bd = Inches(0.56)
    for i, it in enumerate(items):
        r = i // cols; c = i % cols
        x = int(X0 + (cw + gx) * c); y = int(Y0 + (ch + gy) * r); col = PALETTE[i % len(PALETTE)]
        rect(s, x, y, cw, ch, LIGHT); rect(s, x, y, Inches(0.1), ch, col)
        oval(s, x + Inches(0.26), int(y + ch / 2 - bd / 2), bd, bd, col)
        ic = icons[i] if icons else str(i + 1)
        txt(s, x + Inches(0.26), int(y + ch / 2 - bd / 2), bd, bd, [[(ic, 17, WHITE, True)]],
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        tx = x + Inches(1.02); tw = cw - Inches(1.26)
        if isinstance(it, tuple):
            txt(s, tx, int(y + Inches(0.12)), tw, int(ch - Inches(0.18)),
                [[(it[0], size + 1, INK, True)], [(it[1], size - 2, GREY, False)]],
                anchor=MSO_ANCHOR.MIDDLE, space=3)
        else:
            txt(s, tx, int(y + Inches(0.1)), tw, int(ch - Inches(0.16)), [[(it, size, INK, False)]],
                anchor=MSO_ANCHOR.MIDDLE)
    footer(s); return s


def flow_h(title, steps, kicker=None, color=BLUE, note=None):
    """Horizontal numbered process strip — oversized badges so it reads from the back."""
    s = head(slide(), title, kicker, kcolor=color)
    n = len(steps); X0 = Inches(0.85); TOTW = Inches(11.63); gap = Inches(0.34)
    cw = int((TOTW - gap * (n - 1)) / n)
    y = Inches(2.45); ch = Inches(3.15) if not note else Inches(2.85); bd = Inches(0.82)
    for i, st in enumerate(steps):
        x = int(X0 + (cw + gap) * i)
        rect(s, x, y, cw, ch, LIGHT); rect(s, x, y, cw, Inches(0.1), color)
        oval(s, int(x + cw / 2 - bd / 2), int(y + Inches(0.42)), bd, bd, color)
        txt(s, int(x + cw / 2 - bd / 2), int(y + Inches(0.42)), bd, bd, [[(str(i + 1), 30, WHITE, True)]],
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        title_t, sub_t = (st if isinstance(st, tuple) else (st, ""))
        txt(s, x + Inches(0.16), int(y + Inches(1.42)), cw - Inches(0.32), Inches(0.62),
            [[(title_t, 13.5, INK, True)]], align=PP_ALIGN.CENTER)
        if sub_t:
            txt(s, x + Inches(0.16), int(y + Inches(2.02)), cw - Inches(0.32), ch - Inches(2.1),
                [[(sub_t, 11, GREY, False)]], align=PP_ALIGN.CENTER)
        if i < n - 1:
            txt(s, int(x + cw), int(y + Inches(1.15)), gap, Inches(0.6), [[("▶", 15, color, True)]],
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    if note:
        rect(s, Inches(0.85), Inches(5.62), Inches(11.63), Inches(1.05), LIGHT)
        rect(s, Inches(0.85), Inches(5.62), Inches(0.1), Inches(1.05), color)
        txt(s, Inches(1.15), Inches(5.62), Inches(11.2), Inches(1.05), [[(note, 12.5, INK, True)]],
            anchor=MSO_ANCHOR.MIDDLE)
    footer(s); return s


def img_points(title, asset, points, kicker=None, accent=BLUE, caption=None):
    """THE default concept slide: diagram left, takeaway tiles right."""
    s = head(slide(), title, kicker, kcolor=accent)
    p = _asset(asset)
    if p:
        w, h = _fit(p, 7.0, 4.55 if not caption else 4.15)
        s.shapes.add_picture(p, Inches(0.85), Inches(2.0), width=w, height=h)
        if caption:
            rect(s, Inches(0.85), Inches(6.28), Inches(7.0), Inches(0.42), LIGHT)
            txt(s, Inches(1.0), Inches(6.28), Inches(6.7), Inches(0.42), [[(caption, 11, GREY, False)]],
                anchor=MSO_ANCHOR.MIDDLE)
    X = Inches(8.15); W = Inches(4.33)
    n = len(points); gy = Inches(0.2)
    th = int((Inches(4.7) - gy * (n - 1)) / n)
    for i, pt in enumerate(points):
        y = int(Inches(1.98) + (th + gy) * i); col = PALETTE[i % len(PALETTE)]
        rect(s, X, y, W, th, LIGHT); rect(s, X, y, W, Inches(0.09), col)
        t, cap = (pt if isinstance(pt, tuple) else (pt, ""))
        runs = [[(t, 13.5, col, True)]]
        if cap: runs.append([(cap, 11.5, INK, False)])
        txt(s, X + Inches(0.24), y + Inches(0.2), W - Inches(0.46), th - Inches(0.3), runs,
            anchor=MSO_ANCHOR.TOP, space=4)
    footer(s); return s


def img_full(title, asset, kicker=None, accent=BLUE, caption=None):
    s = head(slide(), title, kicker, kcolor=accent)
    p = _asset(asset)
    if p:
        maxh = 4.3 if caption else 4.75
        w, h = _fit(p, 11.6, maxh)
        s.shapes.add_picture(p, int((SW - w) / 2), Inches(1.98), width=w, height=h)
        if caption:
            rect(s, Inches(0.85), Inches(6.36), Inches(11.63), Inches(0.5), LIGHT)
            rect(s, Inches(0.85), Inches(6.36), Inches(0.1), Inches(0.5), accent)
            txt(s, Inches(1.15), Inches(6.36), Inches(11.2), Inches(0.5), [[(caption, 12.5, INK, True)]],
                anchor=MSO_ANCHOR.MIDDLE)
    footer(s); return s


def compare_table(title, headers, rows, kicker=None, accent=BLUE, note=None):
    s = head(slide(), title, kicker, kcolor=accent)
    X = Inches(0.85); W = Inches(11.63)
    ncol = len(headers)
    fracs = [0.26] + [(0.74 / (ncol - 1))] * (ncol - 1)
    widths = [int(W * f) for f in fracs]
    y = Inches(1.98); hh = Inches(0.5)
    rect(s, X, y, W, hh, accent)
    cx = X
    for i, hd in enumerate(headers):
        txt(s, cx + Inches(0.16), y, widths[i] - Inches(0.24), hh, [[(hd, 13, WHITE, True)]],
            anchor=MSO_ANCHOR.MIDDLE)
        cx += widths[i]
    avail = Inches(6.55) - (y + hh) if not note else Inches(5.95) - (y + hh)
    rh = int(min(Inches(0.72), avail / max(1, len(rows))))
    for r, row in enumerate(rows):
        ry = int(y + hh + rh * r)
        rect(s, X, ry, W, rh, LIGHT if r % 2 == 0 else WHITE, line=LINE)
        cx = X
        for i, cell in enumerate(row):
            txt(s, cx + Inches(0.16), ry, widths[i] - Inches(0.24), rh,
                [[(cell, 12, INK, i == 0)]], anchor=MSO_ANCHOR.MIDDLE)
            cx += widths[i]
    if note:
        ny = int(y + hh + rh * len(rows) + Inches(0.18))
        rect(s, X, ny, W, Inches(0.68), LIGHT); rect(s, X, ny, Inches(0.1), Inches(0.68), accent)
        txt(s, X + Inches(0.3), ny, W - Inches(0.5), Inches(0.68), [[(note, 12.5, INK, True)]],
            anchor=MSO_ANCHOR.MIDDLE)
    footer(s); return s


def ncards(title, cards, kicker=None, cols=4, accent=BLUE, synthesis=None):
    """Outlined numbered concept cards — the reference's signature 4-across move."""
    s = head(slide(), title, kicker, kcolor=accent)
    n = len(cards); rows = math.ceil(n / cols)
    X0 = Inches(0.85); TOTW = Inches(11.63); gx = Inches(0.28)
    cw = int((TOTW - gx * (cols - 1)) / cols)
    AREAH = Inches(4.7) if not synthesis else Inches(3.05)
    gy = Inches(0.24); ch = int((AREAH - gy * (rows - 1)) / rows)
    for i, cd in enumerate(cards):
        r = i // cols; c = i % cols
        x = int(X0 + (cw + gx) * c); y = int(Inches(1.98) + (ch + gy) * r)
        col = PALETTE[i % len(PALETTE)]
        rect(s, x, y, cw, ch, WHITE, line=col)
        rect(s, x, y, Inches(0.09), ch, col)
        bd = Inches(0.44)
        oval(s, x + Inches(0.24), y + Inches(0.2), bd, bd, col)
        txt(s, x + Inches(0.24), y + Inches(0.2), bd, bd, [[(str(i + 1), 14, WHITE, True)]],
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        t, body = (cd if isinstance(cd, tuple) else (cd, ""))
        txt(s, x + Inches(0.78), y + Inches(0.16), cw - Inches(0.98), Inches(0.52),
            [[(t, 13, col, True)]], anchor=MSO_ANCHOR.MIDDLE)
        if body:
            txt(s, x + Inches(0.26), y + Inches(0.76), cw - Inches(0.5), ch - Inches(0.9),
                [[(body, 11, INK, False)]])
    if synthesis:
        rect(s, Inches(0.85), Inches(5.32), Inches(11.63), Inches(1.35), LIGHT)
        rect(s, Inches(0.85), Inches(5.32), Inches(0.1), Inches(1.35), accent)
        txt(s, Inches(1.18), Inches(5.46), Inches(11.2), Inches(0.32),
            [[(synthesis[0].upper(), 12, accent, True)]])
        txt(s, Inches(1.18), Inches(5.80), Inches(11.2), Inches(0.8), [[(synthesis[1], 12.5, INK, False)]])
    footer(s); return s


def trainer_slide(kicker, name, role, rows, initials, accent=BLUE):
    s = head(slide(), "About the Trainer", kicker, kcolor=accent)
    lx = Inches(0.85); lw = Inches(3.65)
    rect(s, lx, Inches(1.95), lw, Inches(4.7), LIGHT); rect(s, lx, Inches(1.95), lw, Inches(0.12), accent)
    bd = Inches(1.7); ax = int(lx + (lw - bd) / 2)
    oval(s, ax, Inches(2.5), bd, bd, accent)
    txt(s, ax, Inches(2.5), bd, bd, [[(initials, 42, WHITE, True)]], align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE)
    txt(s, lx + Inches(0.15), Inches(4.55), lw - Inches(0.3), Inches(0.6), [[(name, 20, INK, True)]],
        align=PP_ALIGN.CENTER)
    txt(s, lx + Inches(0.15), Inches(5.2), lw - Inches(0.3), Inches(1.2), [[(role, 13, GREY, False)]],
        align=PP_ALIGN.CENTER)
    rx = Inches(4.9); rw = Inches(7.58); ry = Inches(1.95); rh = Inches(4.7)
    n = len(rows); gy = Inches(0.2); th = int((rh - gy * (n - 1)) / n)
    for i, (label, val) in enumerate(rows):
        y = int(ry + (th + gy) * i); col = PALETTE[i % len(PALETTE)]
        rect(s, rx, y, rw, th, LIGHT); rect(s, rx, y, Inches(0.1), th, col)
        vruns = [(val, 13.5, INK, False)] if val else \
                [("__________________________________________________", 13, LINE, False)]
        txt(s, rx + Inches(0.32), y, rw - Inches(0.6), th,
            [[(label.upper(), 11, col, True)], vruns], anchor=MSO_ANCHOR.MIDDLE, space=3)
    footer(s); return s


def browser_mock(title, url, inner_lines, tiles, kicker=None, accent=BLUE, summary=None):
    """Browser-chrome mock of a web tool/LMS + numbered how-to tiles on the right."""
    s = head(slide(), title, kicker, kcolor=accent)
    BX, BY, BW, BH = Inches(0.85), Inches(1.98), Inches(6.9), Inches(4.3)
    rect(s, BX, BY, BW, BH, WHITE, line=LINE)
    rect(s, BX, BY, BW, Inches(0.46), LIGHT)
    for i, cc in enumerate([RGBColor(0xFF, 0x5F, 0x57), RGBColor(0xFE, 0xBC, 0x2E),
                            RGBColor(0x28, 0xC8, 0x40)]):
        oval(s, BX + Inches(0.16 + i * 0.24), BY + Inches(0.15), Inches(0.16), Inches(0.16), cc)
    roundrect(s, BX + Inches(1.0), BY + Inches(0.1), BW - Inches(1.3), Inches(0.28), WHITE, line=LINE, adj=0.4)
    txt(s, BX + Inches(1.16), BY + Inches(0.1), BW - Inches(1.5), Inches(0.28),
        [[(url, 10.5, BLUE, False)]], anchor=MSO_ANCHOR.MIDDLE)
    iy = BY + Inches(0.72)
    for i, ln in enumerate(inner_lines):
        t, sub = (ln if isinstance(ln, tuple) else (ln, ""))
        col = PALETTE[i % len(PALETTE)]
        rect(s, BX + Inches(0.3), int(iy + Inches(0.78) * i), BW - Inches(0.6), Inches(0.64), LIGHT)
        rect(s, BX + Inches(0.3), int(iy + Inches(0.78) * i), Inches(0.08), Inches(0.64), col)
        runs = [[(t, 12.5, INK, True)]]
        if sub: runs[0].append((f"   {sub}", 10.5, GREY, False))
        txt(s, BX + Inches(0.52), int(iy + Inches(0.78) * i), BW - Inches(0.9), Inches(0.64), runs,
            anchor=MSO_ANCHOR.MIDDLE)
    X = Inches(8.05); W = Inches(4.43)
    n = len(tiles); gy = Inches(0.2); th = int((Inches(4.3) - gy * (n - 1)) / n)
    for i, tl in enumerate(tiles):
        y = int(Inches(1.98) + (th + gy) * i); col = PALETTE[i % len(PALETTE)]
        rect(s, X, y, W, th, LIGHT); rect(s, X, y, Inches(0.09), th, col)
        bd = Inches(0.4)
        oval(s, X + Inches(0.22), int(y + th / 2 - bd / 2), bd, bd, col)
        txt(s, X + Inches(0.22), int(y + th / 2 - bd / 2), bd, bd, [[(str(i + 1), 13, WHITE, True)]],
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        txt(s, X + Inches(0.74), y + Inches(0.12), W - Inches(0.94), th - Inches(0.2),
            [[(tl, 11.5, INK, False)]], anchor=MSO_ANCHOR.MIDDLE)
    if summary:
        rect(s, Inches(0.85), Inches(6.42), Inches(11.63), Inches(0.5), LIGHT)
        rect(s, Inches(0.85), Inches(6.42), Inches(0.1), Inches(0.5), accent)
        txt(s, Inches(1.15), Inches(6.42), Inches(11.2), Inches(0.5), [[(summary, 12, INK, True)]],
            anchor=MSO_ANCHOR.MIDDLE)
    footer(s); return s


def activity_slide(a, topic_code):
    """ONE briefing slide per activity. NO step-by-step — that lives in the Learner Guide."""
    s = head(slide(), a["title"], f"TOPIC {topic_code} · HANDS-ON ACTIVITY", kcolor=TEAL)
    # badge
    rect(s, Inches(9.9), Inches(0.5), Inches(2.58), Inches(0.46), TEAL)
    txt(s, Inches(9.9), Inches(0.5), Inches(2.58), Inches(0.46),
        [[(f"ACTIVITY {a['num']}  ·  {a['duration']} MIN", 11.5, WHITE, True)]],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # scenario — font auto-shrinks to fit; text is NEVER truncated
    rect(s, Inches(0.85), Inches(1.95), Inches(11.63), Inches(1.42), LIGHT)
    rect(s, Inches(0.85), Inches(1.95), Inches(0.1), Inches(1.42), VIOLET)
    txt(s, Inches(1.16), Inches(2.06), Inches(11.2), Inches(0.28), [[("THE SITUATION", 11, VIOLET, True)]])
    txt(s, Inches(1.16), Inches(2.34), Inches(11.15), Inches(0.99),
        [[(a["scenario"], _fit_pt(a["scenario"], 11.15, 0.99, base=12.0, floor=9.0), INK, False)]])
    # what you'll do / produce / tool
    tiles = [(BLUE, "WHAT YOU'LL DO", a["desc"]),
             (TEAL, "YOU'LL PRODUCE", a["build"]),
             (AMBER, "DONE WHEN", a["test"])]
    tw = Inches(3.71); xs = [Inches(0.85), Inches(4.81), Inches(8.77)]
    for (col, lbl, body), x in zip(tiles, xs):
        rect(s, x, Inches(3.55), tw, Inches(2.28), LIGHT); rect(s, x, Inches(3.55), tw, Inches(0.1), col)
        txt(s, x + Inches(0.24), Inches(3.72), tw - Inches(0.45), Inches(0.32), [[(lbl, 11, col, True)]])
        txt(s, x + Inches(0.24), Inches(4.06), tw - Inches(0.45), Inches(1.72),
            [[(body, _fit_pt(body, 3.26, 1.72, base=10.5, floor=8.0), INK, False)]])
    # tool + LG pointer band
    rect(s, Inches(0.85), Inches(6.02), Inches(11.63), Inches(0.78), WHITE, line=TEAL)
    rect(s, Inches(0.92), Inches(6.10), Inches(1.3), Inches(0.62), TEAL)
    txt(s, Inches(0.92), Inches(6.10), Inches(1.3), Inches(0.62),
        [[(f"ACT {a['num']}", 13, WHITE, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Inches(2.35), Inches(6.10), Inches(9.9), Inches(0.34),
        [[(f"Tool: {a['tool']} — {a['tool_url']}", 12, INK, True)]])
    _strap = (f"LO {a['lo'][-1]}  ·  {a['objective']}  ·  "
              f"Full step-by-step: Learner Guide, Activity {a['num']}")
    txt(s, Inches(2.35), Inches(6.42), Inches(9.95), Inches(0.36),
        [[(_strap, _fit_pt(_strap, 9.95, 0.36, base=10.0, floor=7.0), GREY, False)]])
    footer(s); return s


def debrief_slide(a, topic_code):
    """The teaching payoff after each activity — what the activity proves."""
    s = head(slide(), f"Debrief — Activity {a['num']}: What It Proves",
             f"TOPIC {topic_code} · ACTIVITY {a['num']} DEBRIEF", kcolor=VIOLET)
    n = len(a["debrief"])
    gy = Inches(0.26); th = int((Inches(3.4) - gy * (n - 1)) / n)
    for i, d in enumerate(a["debrief"]):
        y = int(Inches(1.98) + (th + gy) * i); col = PALETTE[i % len(PALETTE)]
        rect(s, Inches(0.85), y, Inches(11.63), th, LIGHT)
        rect(s, Inches(0.85), y, Inches(0.1), th, col)
        bd = Inches(0.44)
        oval(s, Inches(1.12), int(y + th / 2 - bd / 2), bd, bd, col)
        txt(s, Inches(1.12), int(y + th / 2 - bd / 2), bd, bd, [[(str(i + 1), 14, WHITE, True)]],
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        txt(s, Inches(1.74), y + Inches(0.1), Inches(10.5), th - Inches(0.2), [[(d, 13, INK, False)]],
            anchor=MSO_ANCHOR.MIDDLE)
    rect(s, Inches(0.85), Inches(5.55), Inches(11.63), Inches(1.18), WHITE, line=GREEN)
    txt(s, Inches(1.15), Inches(5.64), Inches(11.2), Inches(0.3), [[("SUCCESS CRITERION", 11, GREEN, True)]])
    txt(s, Inches(1.15), Inches(5.95), Inches(11.2), Inches(0.72),
        [[(a["test"], _fit_pt(a["test"], 11.2, 0.72, base=11.5, floor=8.5), INK, False)]])
    footer(s); return s


def brk(kind, dur, color=AMBER):
    s = slide(); rect(s, 0, 0, SW, SH, WHITE)
    rect(s, 0, 0, SW, Inches(0.22), color); rect(s, 0, Inches(7.28), SW, Inches(0.22), color)
    rect(s, Inches(5.4), Inches(2.35), Inches(2.53), Inches(0.1), color)
    txt(s, 0, Inches(2.75), SW, Inches(1.2), [[(kind, 46, INK, True)]], align=PP_ALIGN.CENTER)
    txt(s, 0, Inches(4.05), SW, Inches(0.8), [[(dur, 22, color, True)]], align=PP_ALIGN.CENTER)
    PAGE["n"] += 1


# ============================================================ BUILD
cover()

# ---------------------------------------------------------------- ADMIN
mark("admin")
section("COURSE ADMINISTRATION", "Welcome & Housekeeping", "")

tile_grid("Digital Attendance (Mandatory)", [
    ("Three times a day", "Take the AM, PM and Assessment digital attendance — mandatory for every WSQ-funded course."),
    ("Trainer shows the QR", "The trainer or administrator displays the digital attendance QR code from the SSG portal."),
    ("Scan and submit", "Scan the QR code with your mobile phone camera and submit your attendance."),
    ("75% minimum", "A minimum of 75% attendance is required to be eligible for assessment and funding."),
], kicker="ADMIN · TRAQOM & ATTENDANCE", cols=2, size=14)

trainer_slide("ADMIN · YOUR TRAINER (TEMPLATE)", "Your Trainer",
              "WSQ Adult Educator  ·  Agile Practitioner",
              [("Name", ""), ("Certifications", ""), ("Industry experience", ""),
               ("Agile experience", ""), ("Contact", "")], "?", accent=GREY)

trainer_slide("ADMIN · YOUR TRAINER", C.TRAINER,
              "PhD (NUS) · MEng (NTU) · MBA · ACTA/DACE certified adult educator",
              [("Role", "Course developer and principal trainer, Tertiary Infotech Academy"),
               ("Certifications", "Enterprise Design Thinking Practitioner · Agile & Scrum certified · ACTA/DACE"),
               ("Experience", "20+ years across engineering, data analytics, AI and digital transformation"),
               ("Teaches", "Agile project management, design thinking, problem solving and analytics"),
               ("Contact", "enquiry@tertiaryinfotech.com  ·  +65 6100 0613")],
              "AA", accent=BLUE)

tile_grid("Ground Rules", [
    ("Phones on silent", "Set your mobile phone to silent mode for the duration of the session."),
    ("Participate actively", "No question is a stupid question. This course runs on discussion."),
    ("Mutual respect", "Agree to disagree. One conversation at a time."),
    ("Be punctual", "Return from breaks on time so we finish on schedule."),
    ("Step out quietly", "Exit silently for a phone call or a toilet break."),
    ("75% attendance", "Required to be eligible for assessment and funding."),
], kicker="ADMIN · HOW WE WORK TOGETHER", cols=2, size=14)

browser_mock("Download Your Course Material",
             "https://lms-tms.tertiaryinfotech.com",
             [("Log in", "with the email you registered with"),
              ("My Courses", "select Agile Project Management for Business"),
              ("Courseware", "Learner Slides · Learner Guide · Activities"),
              ("Download", "save the PDF to your laptop before we start")],
             ["Open lms-tms.tertiaryinfotech.com in your browser",
              "Log in with your registered email address",
              "Open this course under 'My Courses'",
              "Download the Learner Guide and slides",
              "Keep them open — the assessment is open book"],
             kicker="ADMIN · COURSE MATERIAL", accent=BLUE,
             summary="Everything you need for the activities and the open-book assessment is in the Learner Guide.")

# --- learning outcomes / TSC / outline
tile_grid("Learning Outcomes", [
    ("LO1 — Adopt the Agile mindset", "Adopt a new Agile mindset for project management, and know when Agile fits and when it does not."),
    ("LO2 — Share and implement", "Share and implement Agile practices within your teams — the Manifesto, Scrum, Lean and Kanban."),
    ("LO3 — Build, execute and track", "Build an Agile team to execute and track project performance with real metrics."),
], kicker="BY THE END OF THIS COURSE YOU WILL BE ABLE TO", cols=1, size=15, accent=TEAL)

compare_table("Skills Framework — TSC Abilities Assessed",
              ["Code", "Ability", "Where it is taught"],
              [["A1", "Share information actively within and across teams", "Topic 2 · Activity 4"],
               ["A2", "Organise work in alignment with operational priorities", "Topic 1–2 · Activities 2, 4"],
               ["A3", "Implement Agile or lean practices to reduce waste and defects", "Topic 2–3 · Activities 3, 7"],
               ["A4", "Measure progress against targets on a regular basis", "Topic 3 · Activities 5, 7, 8"],
               ["A5", "Experiment with new ideas, products or services", "Topic 1 · Activity 1"],
               ["A6", "Assess work performance and quality for continuous improvement", "Topic 3 · Activities 6, 7, 8"],
               ["A7", "Manage responsibilities and take ownership of outcomes", "Topic 3 · Activities 5, 6, 8"]],
              kicker=f"TSC: {C.TSC_TITLE}  ·  {C.TSC_CODE}", accent=VIOLET,
              note="Every ability is taught AND practised in a named activity — nothing is assessed that is not taught.")

compare_table("Skills Framework — TSC Knowledge Assessed",
              ["Code", "Knowledge", "Where it is taught"],
              [["K1", "Methods to analyse current and future business operating landscapes", "Topic 1 · Activity 2"],
               ["K2", "Methods to analyse current and future customer needs and preferences", "Topic 1 · Activity 1"],
               ["K3", "Organisational policies, processes and standards", "Topic 1 · Activity 2"],
               ["K4", "Types of change management methodologies, tools and practices", "Topic 2"],
               ["K5", "Types of team composition and formation models", "Topic 3 · Activities 4, 5"],
               ["K6", "Values and principles of Agile methodologies", "Topic 2 · Activity 3"],
               ["K7", "Types of Agile methodologies and practices", "Topic 2 · Activity 3"]],
              kicker=f"TSC: {C.TSC_TITLE}  ·  {C.TSC_CODE}", accent=VIOLET,
              note="Every knowledge code is taught in a topic AND tested by one question in "
                   "the Written Assessment.")

for d in (1, 2):
    # read straight from course_data.DAY_AGENDA — the same single source the Lesson
    # Plan's detailed schedule is built from, so the deck and the LP cannot disagree.
    items = C.DAY_AGENDA[d]
    tile_grid(f"Lesson Plan — Day {d}", items,
              kicker=_ellipsis(f"DAY {d} · {C.DAY_THEMES[d].upper()}", 78), cols=1,
              size=13.5 if len(items) <= 5 else 12.0,
              accent=BLUE if d == 1 else TEAL)

tile_grid("Course Outline", [
    (f"Topic 1 — {C.TOPICS[0]['title']}", C.TOPICS[0]["subtitle"] + f"  ·  {C.TOPICS[0]['tsc']}"),
    (f"Topic 2 — {C.TOPICS[1]['title']}", C.TOPICS[1]["subtitle"] + f"  ·  {C.TOPICS[1]['tsc']}"),
    (f"Topic 3 — {C.TOPICS[2]['title']}", C.TOPICS[2]["subtitle"] + f"  ·  {C.TOPICS[2]['tsc']}"),
], kicker="WHAT WE COVER OVER 2 DAYS · 16 TRAINING HOURS", cols=1, size=14, accent=BLUE)

ncards("The Six Tools You Will Use", [
    (t["name"], f"{t['use']}\n{t['url']}") for t in C.ED_TOOLS
], kicker="HANDS-ON TOOLING · BROWSER-BASED, NOTHING TO INSTALL", cols=3, accent=TEAL,
   synthesis=("HOW WE WORK", "Every activity uses a real browser-based tool on a single running case study — "
                             "HarbourFront Logistics, a Singapore 3PL whose customer portal failed as a waterfall project."))

big_statement("One case study runs through all 8 activities.",
              "HarbourFront Logistics Pte Ltd — a 140-staff Singapore third-party logistics provider whose "
              "'CustomerConnect' portal shipped 4 months late with 62% of its features unused. You will diagnose it, "
              "restart it with Agile, run its sprints, and forecast its release.",
              "THE RUNNING CASE", color=VIOLET)

# --- assessment (briefing BEFORE assessment)
tile_grid("Briefing for Assessment", [
    ("Clear your table", "Place phones and other materials under the table or on the floor."),
    ("No photos or recording", "Assessment scripts must not be photographed or recorded."),
    ("No discussion", "Work individually once the assessment starts."),
    ("Black or blue pen", "Use a black or blue pen for hard-copy scripts."),
    ("No correction fluid", "Do not use liquid paper or correction tape."),
    ("Scripts collected", "Scripts are collected when the time is up."),
], kicker="ASSESSMENT · BRIEFING", cols=2, size=13.5, accent=AMBER)

tile_grid("Final Assessment", [
    ("Written Assessment (WA)", C.ASSESSMENT["written"]),
    ("Case Study (CS)", C.ASSESSMENT["practical"]),
    ("Open book", C.ASSESSMENT["open_book"]),
    ("Competency", C.ASSESSMENT["note"]),
], kicker="ASSESSMENT · WHAT TO EXPECT", cols=1, size=14, accent=AMBER)

img_full("Assessment Flow", "assessment-flow.png", kicker="ASSESSMENT · END-TO-END FLOW",
         accent=AMBER,
         caption="Briefing → Written Assessment → Case Study → marking → result, with an appeal route at every stage.")

tile_grid("Criteria for Funding", [
    ("75% attendance", "Minimum attendance rate of 75%, based on the SSG digital attendance record."),
    ("Assessed Competent", "Complete both assessment instruments and be assessed as 'Competent'."),
    ("TRAQOM survey", "Complete the mandatory course feedback and TRAQOM survey."),
], kicker="ADMIN · FUNDING ELIGIBILITY", cols=1, size=14, accent=GREEN)

# ---------------------------------------------------------------- TOPIC CONCEPT CONTENT
TOPIC_ACTS = {t["num"]: [a for a in ACTIVITIES if a["topic"] == t["num"]] for t in C.TOPICS}


def concept_block(topic_num):
    """The substantive teaching content for a topic — visual grammar, no bullet walls."""
    t = [x for x in C.TOPICS if x["num"] == topic_num][0]
    code = t["code"]

    if topic_num == 1:
        big_statement("Are we building the right thing?",
                      "Traditional project management asks whether we are building the thing right — on time, "
                      "on budget, to specification. Agile asks the prior question first, and keeps asking it.",
                      "TOPIC 01 · THE CENTRAL QUESTION", color=BLUE)
        flow_h("The Traditional Project Lifecycle",
               [("Initiate", "charter, sponsor,\nbusiness case"), ("Plan", "scope, WBS,\nschedule, budget"),
                ("Execute", "build to the\napproved plan"), ("Monitor", "variance against\nthe baseline"),
                ("Close", "handover and\nlessons learned")],
               kicker=f"TOPIC {code} · WHERE WE ARE COMING FROM", color=BLUE,
               note="Nothing here is wrong. It is simply built on an assumption: that the plan written at the "
                    "start remains a good description of reality until the end.")
        tile_grid("The Challenges Every Project Manager Recognises", [
            ("Scope creep", "Requirements change after sign-off, and every change costs a variation order."),
            ("Resource constraints", "Not enough people, or not enough of the right people at the right time."),
            ("Unrealistic schedules", "Dates set by commitment or hope rather than by capacity."),
            ("Unproven technology", "Platforms chosen early and validated far too late."),
            ("Documentation burden", "Reporting effort that grows faster than delivery effort."),
            ("Too many dependencies", "One late input stalls a whole chain of work."),
            ("Stakeholder expectations", "Different stakeholders holding different pictures of 'done'."),
            ("Late discovery", "The biggest problems surface in the final third, when change is most expensive."),
        ], kicker=f"TOPIC {code} · THE PROBLEM SPACE", cols=2, size=13.5, accent=AMBER)
        img_points("Understanding the Waterfall Model", "waterfall-vs-agile.png", [
            ("Scope frozen at the start", "The end product is fixed in all respects when the project begins."),
            ("Detailed upfront estimation", "Requires complete estimation and planning before any build."),
            ("Change is a cost event", "Mid-course changes cause delays and formal variation orders."),
            ("One integration point", "The parts meet late, so integration risk peaks at the worst moment."),
        ], kicker=f"TOPIC {code} · THE WATERFALL MODEL", accent=RED)
        tile_grid("Challenges with the Waterfall Model", [
            ("Change is hostile", "Managing changing requirements mid-lifecycle fights the model rather than using it."),
            ("Late visibility", "Nobody sees the end product until the end of the lifecycle."),
            ("Risk loading", "Risk and uncertainty concentrate at the point where they are most expensive to fix."),
            ("Complexity upfront", "The model asks you to fully understand a complex problem before you have built anything."),
        ], kicker=f"TOPIC {code} · WHY IT STRUGGLES", cols=2, size=14, accent=RED)
        img_full("The Cost of Change", "cost-of-change.png",
                 kicker=f"TOPIC {code} · WHY LATE DISCOVERY IS EXPENSIVE", accent=RED,
                 caption="This single curve is the economic argument for short iterations — it is not about speed, it is about the price of being wrong.")
        ncards("What Agile Actually Is", [
            ("An umbrella term", "Agile covers many iterative, incremental methods — it is not one process."),
            ("Iterative delivery", "APM: 'an iterative approach to delivering a project throughout its life cycle'."),
            ("Built for software, used everywhere", "Now standard in marketing, finance, HR, construction and biotech."),
            ("A mindset, not a ceremony", "Running stand-ups without changing how decisions are made is not Agile."),
        ], kicker=f"TOPIC {code} · DEFINING AGILE", cols=4, accent=TEAL,
           synthesis=("THE CORE IDEA", "Deliver in small increments, get feedback from real users, and let what you "
                                       "learn change what you build next. Everything else in this course is machinery for doing that reliably."))
        img_points("Agile vs Traditional Project Management", "value-delivery.png", [
            ("Builds in increments", "Rather than delivering the whole product as one event."),
            ("Plans continuously", "Planning happens throughout, not only at the start."),
            ("Customer sees value faster", "Benefit is released during the project, not after it."),
            ("Change is welcomed", "Change is expected and priced in, not resisted."),
        ], kicker=f"TOPIC {code} · THE SHIFT", accent=TEAL)
        compare_table("Comparison of Agile and Waterfall",
                      ["Aspect", "Waterfall / Traditional", "Agile"],
                      [["Approach", "Sequential phases, one pass", "Iterative and incremental"],
                       ["Requirements", "Fixed and signed off upfront", "Emerge and are re-ordered continuously"],
                       ["Flexibility", "Resists change after planning", "Adapts throughout"],
                       ["Customer role", "Consulted at start and at UAT", "Continuous collaboration"],
                       ["Delivery", "One release at the end", "A usable increment every sprint"],
                       ["Risk", "Addressed upfront, realised late", "Retired continuously"],
                       ["Documentation", "Comprehensive by mandate", "Barely sufficient, just in time"],
                       ["Team structure", "Hierarchical, specialised roles", "Self-organising, cross-functional"]],
                      kicker=f"TOPIC {code} · SIDE BY SIDE", accent=BLUE,
                      note="Neither column is 'better'. The right question is which column matches the uncertainty of YOUR work.")
        img_points("Inverting the Triangle", "iron-triangle.png", [
            ("Waterfall fixes scope", "Time and cost flex to deliver the agreed feature list."),
            ("Agile fixes time and cost", "You buy a stable team and a timebox; scope is the variable."),
            ("This is the real trade", "You cannot fix all three and still respond to what you learn."),
            ("Implication for governance", "Funding shifts from one big approval to incremental funding decisions."),
        ], kicker=f"TOPIC {code} · THE PARADIGM SHIFT", accent=VIOLET)
        tile_grid("Analysing the Business Operating Landscape (K1)", [
            ("PESTLE", "Political, economic, social, technological, legal and environmental forces reshaping the market."),
            ("SWOT", "Internal strengths and weaknesses against external opportunities and threats."),
            ("Porter's Five Forces", "Competitive rivalry, supplier and buyer power, substitutes and new entrants."),
            ("Scenario planning", "Several plausible futures rather than one forecast, each with a response."),
            ("VUCA framing", "Volatility, uncertainty, complexity and ambiguity — name which one you actually face."),
            ("Re-run the analysis", "The Agile difference: quarterly, not once per project."),
        ], kicker=f"TOPIC {code} · K1 · METHODS AND TOOLS", cols=2, size=13, accent=BLUE)
        tile_grid("Analysing Customer Needs and Preferences (K2)", [
            ("Personas & empathy maps", "What the customer says, thinks, does and feels — used in Activity 1."),
            ("Journey maps", "The end-to-end experience, exposing where the pain actually sits."),
            ("Jobs to be done", "The progress the customer is trying to make, not the feature they requested."),
            ("Kano analysis", "Sorting delighters, satisfiers, must-haves and indifferent features."),
            ("A/B experiments", "Testing a preference with real behaviour instead of asking an opinion."),
            ("Continuous feedback", "Sprint reviews turn customer contact into a rhythm, not an event."),
        ], kicker=f"TOPIC {code} · K2 · METHODS AND TOOLS", cols=2, size=13, accent=TEAL)
        tile_grid("Agile Beyond Software — Business Use Cases", [
            ("Adobe · marketing", "Campaign teams run sprints and retrospectives to adapt messaging to live data."),
            ("Toyota · manufacturing", "The origin of Lean and Kanban — flow, WIP limits and waste elimination."),
            ("Spotify · product delivery", "Autonomous squads optimised for delivery speed and learning."),
            ("PayPal · workforce alignment", "Agile used to align large distributed teams on shared outcomes."),
            ("Banking · compliance delivery", "Hybrid Agile inside regulated stage-gate governance."),
            ("Construction & biotech", "Iterative design and staged validation where rework is costly."),
        ], kicker=f"TOPIC {code} · IT IS NOT JUST FOR SOFTWARE", cols=2, size=13, accent=VIOLET)
        tile_grid("The Agile Mindset", [
            "Welcoming change rather than resisting it", "Working in small increments of real value",
            "Using build-and-feedback loops to learn", "Learning through discovery instead of assumption",
            "Value-driven development, not activity-driven", "Failing fast, and extracting the lesson",
            "Continuous delivery of working outcomes", "Continuous improvement of how the team works",
        ], kicker=f"TOPIC {code} · HOW AGILE PEOPLE THINK", cols=2, size=14, accent=TEAL)
        ncards("When Agile Is NOT the Right Answer", [
            ("Truly fixed scope", "A regulatory submission with a legislated, unchangeable specification."),
            ("No customer access", "If nobody can give feedback each sprint, the feedback loop is theatre."),
            ("Physical irreversibility", "Where an increment cannot be built and changed cheaply."),
            ("No mandate to change governance", "Agile teams inside stage-gate funding produce reports, not agility."),
        ], kicker=f"TOPIC {code} · INTELLECTUAL HONESTY", cols=4, accent=RED,
           synthesis=("CHOOSE DELIBERATELY", "Stable requirements and a documentation mandate favour waterfall; evolving "
                                             "requirements and high novelty favour Agile. Most real organisations run a considered hybrid."))

    elif topic_num == 2:
        big_statement("Four values. Twelve principles. Everything else is machinery.",
                      "The Agile Manifesto was written in 2001 by 17 practitioners. It is one page long, and it is "
                      "the reference every framework in this topic is derived from.",
                      "TOPIC 02 · THE FOUNDATION", color=TEAL)
        img_full("The Agile Manifesto — Four Values", "manifesto-values.png",
                 kicker=f"TOPIC {code} · K6 · THE FOUR VALUES", accent=BLUE,
                 caption="agilemanifesto.org — 'while there is value in the items on the right, we value the items on the left more'.")
        ncards("Reading the Four Values Correctly", [
            ("Individuals & interactions", "Tools are necessary; they never rescue a team that will not talk. Problems are solved by people."),
            ("Working software", "Deliver something that works. Keep documents barely sufficient, just in time, and for a stated reason."),
            ("Customer collaboration", "Manage change rather than suppress it. Build a shared, written definition of 'done'."),
            ("Responding to change", "Energy spent defending the original plan is energy taken from delivering value."),
        ], kicker=f"TOPIC {code} · WHAT EACH VALUE MEANS IN PRACTICE", cols=4, accent=BLUE,
           synthesis=("THE COMMON MISREADING", "'We value the left over the right' does NOT mean the right is worthless. "
                                               "Agile teams still plan, still document, still contract — they just refuse to let those things outrank delivered value."))
        tile_grid("The Twelve Principles — 1 to 6", [
            ("1 · Satisfy the customer", "Through early and continuous delivery of valuable work."),
            ("2 · Welcome change", "Even late in development — harness change for competitive advantage."),
            ("3 · Deliver frequently", "From a couple of weeks to a couple of months; prefer the shorter timescale."),
            ("4 · Work together daily", "Business people and delivery people, throughout the project."),
            ("5 · Motivated individuals", "Give them the environment and support they need, and trust them."),
            ("6 · Face-to-face conversation", "The most efficient and effective way to convey information."),
        ], kicker=f"TOPIC {code} · K6 · PRINCIPLES 1–6", cols=2, size=13, accent=TEAL)
        tile_grid("The Twelve Principles — 7 to 12", [
            ("7 · Working output measures progress", "Not percent-complete of tasks, and not documents produced."),
            ("8 · Sustainable pace", "Sponsors, developers and users maintain a constant pace indefinitely."),
            ("9 · Technical excellence", "Continuous attention to quality and good design enhances agility."),
            ("10 · Simplicity", "Maximising the work NOT done is essential."),
            ("11 · Self-organising teams", "The best architectures, requirements and designs emerge from them."),
            ("12 · Reflect and adjust", "At regular intervals, tune and adjust behaviour accordingly."),
        ], kicker=f"TOPIC {code} · K6 · PRINCIPLES 7–12", cols=2, size=13, accent=TEAL)
        img_full("The Scrum Framework", "scrum-framework.png",
                 kicker=f"TOPIC {code} · K7 · SCRUM END TO END", accent=VIOLET,
                 caption="Scrum is built on empiricism: transparency, inspection and adaptation. The loop is the framework.")
        img_points("Why We Teach Scrum in Depth", "framework-adoption.png", [
            ("Scrum dominates adoption", "Roughly 63% of Agile teams use Scrum, so it is the shared vocabulary."),
            ("It is prescriptive enough to learn", "Fixed roles, artefacts and events give a new team something to follow."),
            ("It generalises", "The same cadence works for marketing, operations and product teams."),
            ("Then tailor it", "Shu-Ha-Ri: follow the framework first, adapt once you understand why."),
        ], kicker=f"TOPIC {code} · FRAMEWORK CHOICE", accent=BLUE)
        ncards("The Three Scrum Accountabilities", [
            ("Product Owner", "Owns value and the ORDER of the product backlog. One person, not a committee. Accepts or rejects each increment."),
            ("Scrum Master", "Owns the process. A servant leader who removes impediments, coaches the team and shields it from interruption."),
            ("Developers", "Own HOW the work is done and how much enters the sprint. Cross-functional, self-organising, collectively accountable."),
        ], kicker=f"TOPIC {code} · WHO DOES WHAT", cols=3, accent=VIOLET,
           synthesis=("THE FAILURE MODE", "Almost every 'Scrum doesn't work here' story traces to a Product Owner who cannot "
                                          "decide, or a manager who overrules the sprint. You will map this yourself in Activity 4."))
        compare_table("Scrum Artefacts and Their Commitments",
                      ["Artefact", "What it is", "Its commitment"],
                      [["Product Backlog", "The ordered list of everything known to be needed", "The Product Goal"],
                       ["Sprint Backlog", "The selected items plus the plan to deliver them", "The Sprint Goal"],
                       ["Increment", "A usable, integrated slice of the product", "The Definition of Done"]],
                      kicker=f"TOPIC {code} · THE THREE ARTEFACTS", accent=TEAL,
                      note="Each artefact carries a commitment. An artefact without its commitment becomes a list nobody trusts.")
        flow_h("The Five Scrum Events",
               [("Sprint Planning", "what and how\n· up to 8h"), ("Daily Scrum", "15 min · re-plan\nthe next 24h"),
                ("The Work", "build to the\nDefinition of Done"), ("Sprint Review", "inspect the product\nwith stakeholders"),
                ("Retrospective", "inspect the\nprocess")],
               kicker=f"TOPIC {code} · THE SPRINT CONTAINER", color=VIOLET,
               note="All five sit inside the Sprint — a fixed-length container of 1–4 weeks. The length does not change sprint to sprint.")
        img_full("Definition of Done", "definition-of-done.png",
                 kicker=f"TOPIC {code} · THE QUALITY BAR", accent=GREEN,
                 caption="The DoD is the team's shared quality bar. Without it, 'done' means something different to every person in the room.")
        tile_grid("Lean — Seven Principles from Toyota", [
            ("Eliminate waste", "To maximise value, minimise everything that is not value."),
            ("Amplify learning", "Communicate early and often; get feedback as soon as possible."),
            ("Decide as late as possible", "Keep options open until the last responsible moment."),
            ("Deliver as fast as possible", "Speed and quality are allies, not opposites."),
            ("Empower the team", "Respect the team's superior knowledge of the technical steps."),
            ("Build quality in", "Assure quality throughout, rather than inspecting it at the end."),
            ("Optimise the whole", "The system is more than the sum of its parts — beware local optimisation."),
            ("Kaizen", "Small, frequent, team-owned improvement, forever."),
        ], kicker=f"TOPIC {code} · K7 · LEAN", cols=2, size=13, accent=AMBER)
        img_full("The Eight Wastes", "lean-wastes.png",
                 kicker=f"TOPIC {code} · WHAT LEAN ATTACKS", accent=AMBER,
                 caption="In knowledge work the dominant waste is almost always WAITING — which is why value stream mapping is so effective.")
        img_full("Kanban — Visualise, Limit WIP, Manage Flow", "kanban-littles-law.png",
                 kicker=f"TOPIC {code} · K7 · KANBAN", accent=BLUE,
                 caption="Little's Law is the engine: limiting WIP is what actually makes delivery faster, not working harder.")
        ncards("Kanban's Five Practices", [
            ("Visualise the workflow", "Knowledge work is invisible until you put it on a board."),
            ("Limit work in progress", "WIP limits surface bottlenecks instead of hiding them."),
            ("Manage flow", "Track how work moves, and measure the effect of each change."),
            ("Make policies explicit", "Write down what 'ready' and 'done' mean so the team can debate them."),
            ("Improve collaboratively", "Evolve the process experimentally, as a team."),
        ], kicker=f"TOPIC {code} · HOW KANBAN WORKS", cols=5, accent=BLUE,
           synthesis=("THE COUNTER-INTUITIVE PART", "Starting less work finishes more work. You will prove this to yourself "
                                                    "in Activity 5 when the burndown flattens, and again in Activity 8 with the control chart."))
        compare_table("Scrum vs Kanban vs Scrumban — Choosing",
                      ["Dimension", "Scrum", "Kanban", "Scrumban"],
                      [["Cadence", "Fixed sprints, 1–4 weeks", "Continuous flow", "Cadence plus WIP limits"],
                       ["Commitment", "A Sprint Goal per sprint", "Pull when capacity frees", "Goal, with flexible pull"],
                       ["Best for", "Teams of 5–9 with a product goal", "Support queues, high variability", "Teams outgrowing strict Scrum"],
                       ["Roles", "PO, SM, Developers", "No prescribed roles", "Usually keeps PO and SM"],
                       ["Key metric", "Velocity, sprint burndown", "Cycle time, throughput, CFD", "Both sets"]],
                      kicker=f"TOPIC {code} · K7 · FRAMEWORK SELECTION", accent=TEAL,
                      note="Pick by the shape of the work, never by fashion. A support team forced into sprints will fight the framework every week.")
        tile_grid("Extreme Programming (XP)", [
            ("Five values", "Simplicity, communication, feedback, courage and respect."),
            ("Test-driven development", "Write the test first; the code passes it once written correctly."),
            ("Pair programming", "Two people, real-time review, and knowledge spread through the team."),
            ("Continuous integration", "Bring code together constantly so incompatibility surfaces early."),
            ("Refactoring", "Remove redundancy and rejuvenate design continuously, not 'later'."),
            ("Collective ownership", "Anyone can improve any part; less risk when someone leaves."),
            ("Small releases", "Frequent small releases keep progress visible to the customer."),
            ("Sustainable pace", "Repeated long hours are unsustainable and counterproductive."),
        ], kicker=f"TOPIC {code} · K7 · XP", cols=2, size=13, accent=VIOLET)
        ncards("Other Methods Worth Knowing", [
            ("DSDM / AgilePM", "The APM-endorsed framework: feasibility, foundations, evolutionary development, deployment. Strong on governance."),
            ("Feature-Driven Development", "Model first, build a feature list, then plan and build by feature."),
            ("Crystal", "A family of methods sized by team and criticality, colour-coded."),
            ("SAFe / LeSS", "Scaling frameworks for many teams on one product or portfolio."),
        ], kicker=f"TOPIC {code} · K7 · THE WIDER LANDSCAPE", cols=4, accent=GREY,
           synthesis=("FOR SINGAPORE PRACTITIONERS", "DSDM/AgilePM is the framework APM (UK) endorses and the one most often "
                                                     "referenced in regulated and public-sector delivery — worth knowing by name in a tender conversation."))
        # change management / resistance (K4)
        tile_grid("Finding Agile Support — Address Each Fear Directly (K4)", [
            ("Senior management & sponsors", "Fear: the risk of visible failure. Answer: incremental funding lets you stop early and cheaply."),
            ("Middle managers", "Fear: loss of control. Answer: they gain real progress data instead of percent-complete estimates."),
            ("The project team", "Fear: exposure and surveillance. Answer: the board shows the WORK, not individual performance."),
            ("Users & customers", "Fear: losing promised features. Answer: they get the highest-value features sooner, and a say each sprint."),
            ("Finance & PMO", "Fear: no fixed scope to audit. Answer: fixed cadence, fixed cost, auditable increments."),
            ("Quality & compliance", "Fear: reduced rigour. Answer: the Definition of Done makes the quality bar explicit and testable."),
        ], kicker=f"TOPIC {code} · K4 · PREPARING FOR RESISTANCE", cols=2, size=12.5, accent=AMBER)
        ncards("Change Management — Making the Adoption Stick (K4)", [
            ("Start with a willing pilot", "Volunteers, on a real project that matters, with visible sponsorship."),
            ("Make the wins visible", "Publish the increment and the metrics. Evidence beats evangelism."),
            ("Evolutionary, not revolutionary", "APM's guidance: change the culture in steps the organisation can absorb."),
            ("Train, coach, then withdraw", "Coach the team through 2–3 sprints, then let it own the process."),
        ], kicker=f"TOPIC {code} · K4 · METHODOLOGIES AND PRACTICES", cols=4, accent=VIOLET,
           synthesis=("THE HONEST WARNING", "Never sell Agile as a way to get more output for less money. It buys "
                                            "responsiveness, earlier value and lower risk of building the wrong thing — not free capacity."))

    else:
        big_statement("A team, a vision, a rhythm, and honest numbers.",
                      "Topic 3 is where Agile stops being a philosophy and becomes a week of work: who is on the team, "
                      "what they are aiming at, how the work flows, and what the metrics actually tell you.",
                      "TOPIC 03 · MAKING IT RUN", color=VIOLET)
        ncards("Team Composition and Formation Models (K5)", [
            ("Cross-functional", "Every skill needed to get to Done sits inside the team."),
            ("Self-organising", "The team decides how to do the work, not just what to do."),
            ("Fewer than ~12", "Communication paths grow faster than headcount. Small teams stay coherent."),
            ("Stable membership", "Velocity is meaningless if the team changes every sprint."),
            ("Generalising specialists", "Deep in one skill, capable across several — this removes queues."),
            ("Co-located or deliberately connected", "Osmotic communication, or tooling that replaces it on purpose."),
        ], kicker=f"TOPIC {code} · K5 · WHAT A GOOD AGILE TEAM LOOKS LIKE", cols=3, accent=TEAL,
           synthesis=("THE BOTTLENECK YOU WILL MEET", "One specialist per skill guarantees queues — work piles up behind "
                                                      "the only person who can do it. You will see exactly this as the Testing bottleneck in Activity 8."))
        img_full("Tuckman's Stages and Adaptive Leadership", "tuckman-leadership.png",
                 kicker=f"TOPIC {code} · K5 · TEAM DEVELOPMENT", accent=BLUE,
                 caption="Match your leadership style to the team's stage. Delegating to a forming team fails; directing a performing team insults it.")
        tile_grid("Servant Leadership — What the Agile Leader Actually Does", [
            ("Shield the team", "Absorb interruptions so the team can hold a sustainable focus."),
            ("Remove impediments", "Own the blockers the team cannot clear itself, and clear them fast."),
            ("Re-communicate the vision", "People forget the goal under delivery pressure. Repeat it."),
            ("Tap intrinsic motivation", "Autonomy, mastery and purpose outperform instruction."),
            ("Model the behaviour", "Honesty, competence and forward-looking judgement, visibly."),
            ("Create safety to experiment", "A team that cannot fail safely will not tell you the truth."),
        ], kicker=f"TOPIC {code} · LEADING WITHOUT AUTHORITY", cols=2, size=13, accent=VIOLET)
        tile_grid("Setting the Shared Vision", [
            ("Agile charter", "Who, what, why, when, where and how — high-level, and authority to proceed."),
            ("Product vision statement", "One paragraph any team member can repeat from memory."),
            ("Definition of Done", "The shared, testable meaning of 'finished'."),
            ("Personas", "Grounded, goal-oriented descriptions of the real users."),
            ("Wireframes & prototypes", "Low-fidelity artefacts that make 'done' visible before it is built."),
            ("Journey and story maps", "The whole experience, so the team can see what each slice contributes."),
        ], kicker=f"TOPIC {code} · GETTING ONE PICTURE IN EVERY HEAD", cols=2, size=13, accent=BLUE)
        img_full("The Anatomy of a User Story", "user-story-anatomy.png",
                 kicker=f"TOPIC {code} · REQUIREMENTS AS STORIES", accent=BLUE,
                 caption="Role, goal, benefit — plus testable acceptance criteria. This exact story is the one you write in Activity 3.")
        ncards("The Three Cs and INVEST", [
            ("Card", "The story is a placeholder, deliberately too small to hold the full requirement."),
            ("Conversation", "The detail is developed by talking, not by writing a longer document."),
            ("Confirmation", "Acceptance criteria state how you will know it is done."),
            ("INVEST", "Independent · Negotiable · Valuable · Estimatable · Small · Testable."),
        ], kicker=f"TOPIC {code} · WRITING STORIES THAT WORK", cols=4, accent=TEAL,
           synthesis=("THE SPLIT RULE", "A story too large for one sprint must be split by WORKFLOW STEP or by DATA TYPE — "
                                        "never by technical layer. Splitting into 'the database part' and 'the UI part' produces a sprint with nothing to demonstrate."))
        tile_grid("Prioritisation Techniques", [
            ("MoSCoW", "Must / Should / Could / Won't have this time. Used in Activity 3."),
            ("Dot voting", "Each person distributes a fixed number of dots across the options."),
            ("100-point method", "Each stakeholder allocates 100 points, forcing real trade-offs."),
            ("Monopoly money", "Equal 'funds' spent on what each stakeholder values most."),
            ("Kano analysis", "Separates delighters, satisfiers, must-haves and indifferent features."),
            ("Weighted shortest job first", "Value divided by effort — highest ratio goes first."),
        ], kicker=f"TOPIC {code} · ORDERING THE BACKLOG", cols=2, size=13, accent=AMBER)
        ncards("Estimation — Why Relative Beats Absolute", [
            ("Humans estimate badly in absolutes", "We are poor at hours, and consistently good at comparisons."),
            ("Story points", "Complexity, effort and risk in one relative number."),
            ("Planning poker", "Simultaneous reveal on a Fibonacci scale defeats anchoring and the loudest voice."),
            ("Affinity & T-shirt sizing", "Fast grouping for large backlogs before detailed estimation."),
        ], kicker=f"TOPIC {code} · AGILE ESTIMATION", cols=4, accent=VIOLET,
           synthesis=("THE 13-POINT RULE", "A story estimated at 13 or 21 points is not a big story — it is a story the team "
                                           "does not yet understand. Split it before it enters a sprint."))
        img_points("Timeboxing and Parkinson's Law", "retro-stages.png", [
            ("Daily Scrum · 15 minutes", "Long enough to re-plan a day, short enough to stay standing."),
            ("Retrospective · ~2 hours", "For a two-week sprint, in five deliberate stages."),
            ("Sprint · 1–4 weeks", "Fixed length, so velocity means something across sprints."),
            ("Parkinson's Law", "Work expands to fill the time available — so the box does the managing."),
        ], kicker=f"TOPIC {code} · THE DISCIPLINE OF THE TIMEBOX", accent=AMBER)
        big_statement("Five metrics tell you everything you need to know.",
                      "Sprint burndown for the sprint. Release burndown for the release. Velocity for the forecast. "
                      "Control chart for cycle time. Cumulative flow diagram for the bottleneck.",
                      f"TOPIC {code} · A4 · MEASURING PROGRESS", color=BLUE)
        img_points("Metric 1 — Sprint Burndown", "sprint-burndown.png", [
            ("What it shows", "Points remaining inside one sprint, day by day, against the ideal line."),
            ("Read the flat line", "Flat means work is STARTED but not FINISHED — check WIP, not effort."),
            ("Read the cliff", "A late vertical drop means work was integrated at the end. That is a risk, not a triumph."),
            ("Carryover", "Unfinished work returns to the backlog at full estimate. It is never partial velocity."),
        ], kicker=f"TOPIC {code} · A4 · TRACKING THE SPRINT", accent=BLUE)
        img_points("Metrics 2 & 3 — Velocity and the Release Forecast", "velocity-forecast.png", [
            ("Velocity", "Average points DONE per sprint. A forecasting input, never a performance target."),
            ("Forecast as a range", "Slowest, average and fastest sprints give a cone, not a date."),
            ("State the assumptions", "Stable team, stable backlog, same estimation scale, no lost sprints."),
            ("The anti-pattern", "Reward velocity and you get estimate inflation — the forecast then breaks."),
        ], kicker=f"TOPIC {code} · A4 · FORECASTING HONESTLY", accent=TEAL)
        img_points("Metric 4 — The Control Chart", "control-chart.png", [
            ("What it shows", "Cycle time and lead time per work item, and the trend across sprints."),
            ("Lead vs cycle time", "Lead time is the customer's whole wait; cycle time is the team's active portion."),
            ("The warning sign", "Velocity rising AND cycle time rising means WIP is rising."),
            ("What to do", "Lower the WIP limit at the constraint. Stop starting, start finishing."),
        ], kicker=f"TOPIC {code} · A4 · CYCLE TIME", accent=VIOLET)
        img_points("Metric 5 — The Cumulative Flow Diagram", "cumulative-flow.png", [
            ("What it shows", "The count of items in each status over time, as stacked bands."),
            ("Read the widening band", "A band widening vertically is a bottleneck at THAT status."),
            ("Read the flat Done band", "Flat 'Done' means nothing is being completed, whatever the team is busy with."),
            ("Act on the constraint", "Move capacity to the widening column — not to the busiest-looking one."),
        ], kicker=f"TOPIC {code} · A4 · FINDING THE BOTTLENECK", accent=AMBER)
        compare_table("Little's Law — Why Limiting WIP Speeds Delivery",
                      ["Scenario", "WIP", "Throughput", "Cycle time"],
                      [["Team starts everything", "30 items", "5 items/week", "6 weeks"],
                       ["Team halves WIP", "15 items", "5 items/week", "3 weeks"],
                       ["Team quarters WIP", "8 items", "5 items/week", "1.6 weeks"]],
                      kicker=f"TOPIC {code} · CYCLE TIME = WIP ÷ THROUGHPUT", accent=BLUE,
                      note="Throughput did not change in any row. Only WIP changed — and cycle time fell with it. This is arithmetic, not opinion.")
        img_full("The Retrospective in Five Stages", "retro-stages.png",
                 kicker=f"TOPIC {code} · A6 · CONTINUOUS IMPROVEMENT", accent=GREEN,
                 caption="The retrospective is the engine of improvement. Without an owned, funded action, it is just a feelings meeting.")
        ncards("Root-Cause Tools — Use Them Together", [
            ("5 Whys", "Drills ONE causal chain to a process cause you can change. Activity 6."),
            ("Fishbone", "Spreads causes across categories to see the whole problem. Activity 2."),
            ("Pareto", "Shows which few causes carry most of the pain. Activity 7."),
            ("PDCA / Kaizen", "Plan-Do-Check-Act: the loop that carries the fix into the next sprint."),
        ], kicker=f"TOPIC {code} · A6 · DIAGNOSIS BEFORE ACTION", cols=4, accent=VIOLET,
           synthesis=("THE SEQUENCE THAT WORKS", "Fishbone to see the whole space, Pareto to pick the few that matter, "
                                                 "5 Whys to reach the changeable cause, then one SMART action with an owner and points in the next sprint."))
        img_points("Pareto — Which Few Causes Carry Most of the Pain", "pareto-defects.png", [
            ("The 80/20 pattern", "A small number of causes generate most of the defects — reliably, across domains."),
            ("Always sort descending", "An unsorted Pareto chart cannot be read. The cumulative line depends on the order."),
            ("Find the crossing point", "Read where the cumulative line crosses 80% — those are the vital few."),
            ("Decide, don't admire", "The chart's purpose is to justify what you will NOT fix this sprint."),
        ], kicker=f"TOPIC {code} · A3 · PRIORITISING THE FIX", accent=RED,
            caption="The real CustomerConnect defect data you will analyse in Activity 7 — 120 defects across 8 causes.")
        img_full("Value Stream Mapping — Where the Time Actually Goes", "value-stream-map.png",
                 kicker=f"TOPIC {code} · A3 · REDUCING WASTE", accent=RED,
                 caption="13 days of work inside 125 days of lead time — 10% efficiency. The opportunity is in the waiting, not the working.")
        tile_grid("Risk as Anti-Value", [
            ("Risk is negative value", "It erodes or removes value, so it belongs in the same conversation as value."),
            ("Risk-adjusted backlog", "High-risk, high-value items are pulled earlier to retire uncertainty."),
            ("Expected monetary value", "Probability × impact, so risks can be compared honestly."),
            ("Risk burndown chart", "Tracks total exposure falling sprint by sprint."),
            ("Spikes", "A timeboxed investigation to reduce a specific technical or risk unknown."),
            ("Schedule risk work as real items", "A risk register nobody funds changes nothing."),
        ], kicker=f"TOPIC {code} · MANAGING UNCERTAINTY", cols=2, size=13, accent=RED)
        ncards("Technical Debt and Ownership (A7)", [
            ("Debt compounds", "Work skipped to go faster makes every later sprint slower."),
            ("Fund refactoring inside the sprint", "'Later' never arrives on its own."),
            ("Collective commitment", "The team commits to the Sprint Goal together, then each member pulls their own work."),
            ("Raise impediments same-day", "Surfacing a blocker on day 9 is the failure — not the blocker itself."),
        ], kicker=f"TOPIC {code} · A7 · TAKING OWNERSHIP", cols=4, accent=TEAL,
           synthesis=("WHAT ACCOUNTABILITY LOOKS LIKE", "An honest board, a same-day impediment, and a 'Done' that means Done. "
                                                        "Every metric in this topic depends on those three behaviours."))


# ---------------------------------------------------------------- TOPICS + ACTIVITIES
for ti, t in enumerate(C.TOPICS):
    mark(f"topic{t['num']}")
    section(f"TOPIC {t['code']}", t["title"], t["code"],
            t["subtitle"] + f"   ·   {t['lo']}   ·   {t['tsc']}")
    concept_block(t["num"])

    acts = TOPIC_ACTS[t["num"]]
    if acts:
        # ncards adapts to the real number of activities — never pad with an empty card
        ncards(f"Hands-On Activities — {t['title']}",
               [(f"Activity {a['num']} · {a['title']}",
                 f"Tool: {a['tool']}\n{a['duration']} minutes  ·  {a['lo']}") for a in acts],
               kicker="WHAT YOU'LL DO", cols=min(len(acts), 2) if len(acts) <= 2 else 2,
               accent=TEAL,
               synthesis=("ALL ON ONE CASE", "Every activity in this topic advances the same HarbourFront "
                                             "CustomerConnect case, so the artefacts you build compound."))
        for a in acts:
            mark(f"act{a['num']}")
            # the ed-tool: the REAL screenshot of the tool, with how-to tiles beside it
            img_points(f"The Tool — {a['tool']}", TOOL_SHOT[a["tool"]], [
                ("Open it in your browser", a["tool_url"]),
                ("No install, no login", "It runs in the browser. Work on one shared screen as a team."),
                ("Build the artefact", "The deliverable named on the next slide."),
                ("Export it", "Screenshot or export into your Activity worksheet."),
            ], kicker=f"ACTIVITY {a['num']} · TOOL", accent=TEAL,
                caption=f"Activity {a['num']} · {a['duration']} minutes · "
                        f"full step-by-step: Learner Guide, Activity {a['num']}")
            activity_slide(a, t["code"])
            debrief_slide(a, t["code"])

    # topic recap — one tile per activity, deduplicated, with the ability codes it evidences
    tile_grid(f"Recap — {t['title']}",
              [(f"Activity {a['num']} · {a['title']}", a["objective"]) for a in acts],
              kicker=f"TOPIC RECAP  ·  {t['lo']}  ·  {t['tsc']}", cols=1,
              size=12.5 if len(acts) > 3 else 14,
              accent=BLUE)

    # breaks
    if t["num"] == 1:
        brk("Lunch Break", "1 hour  ·  Digital attendance (PM) when you return", color=AMBER)
    elif t["num"] == 2:
        brk("End of Day 1", "See you at 09:30 tomorrow  ·  Digital attendance (AM)", color=BLUE)

# ---------------------------------------------------------------- CLOSE
mark("close")
section("WRAP-UP", "Course Summary & Next Steps", "")

tile_grid("What You Can Now Do", [
    ("LO1 — Adopt the Agile mindset", "You can explain why waterfall concentrates risk, diagnose a failed delivery, and reframe requirements around customer value."),
    ("LO2 — Share and implement practices", "You can apply the Manifesto, run Scrum, build an ordered backlog, plan a sprint and clarify accountability."),
    ("LO3 — Build, execute and track", "You can run a sprint, read all five Agile metrics, find a root cause and forecast a release honestly."),
], kicker="LEARNING OUTCOMES ACHIEVED", cols=1, size=14, accent=TEAL)

ncards("The Eight Things Worth Remembering", [
    ("Fixed time, variable scope", "That inversion is the whole paradigm shift."),
    ("Value early beats value complete", "Shipping 20% of scope in month two changes the return profile."),
    ("The sprint goal makes a sprint", "Without one you have a task list with a deadline."),
    ("Done means Done", "Carryover is never velocity, and 'nearly' is not a status."),
    ("One Accountable per decision", "Two Product Owners is the most common Agile failure."),
    ("Limit WIP to go faster", "Little's Law: cycle time = WIP ÷ throughput."),
    ("Velocity forecasts, never targets", "Reward it and you destroy it."),
    ("Actions need owners", "An improvement with no owner and no capacity is a wish."),
], kicker="THE TAKEAWAYS", cols=4, accent=BLUE,
   synthesis=("IF YOU REMEMBER ONE THING", "Agile is not a faster way to build what you were told to build. It is a "
                                           "disciplined way to keep finding out what is actually worth building — and to stop building what is not."))

tile_grid("Where to Go Next", [
    (rc, "www.tertiarycourses.com.sg") for rc in C.RECOMMENDED
], kicker="RECOMMENDED COURSES", cols=1, size=12.5, accent=VIOLET)

# --- closing admin block: Assessment → Assessment Flow → TRAQOM → Thank You
tile_grid("Final Assessment", [
    ("Written Assessment (WA)", C.ASSESSMENT["written"]),
    ("Case Study (CS)", C.ASSESSMENT["practical"]),
    ("Open book", C.ASSESSMENT["open_book"]),
    ("Competency", C.ASSESSMENT["note"]),
], kicker="ASSESSMENT · NOW", cols=1, size=14, accent=AMBER)

img_full("Assessment Flow", "assessment-flow.png", kicker="ASSESSMENT · END-TO-END FLOW",
         accent=AMBER,
         caption="You must be assessed Competent in BOTH instruments. Feedback and an appeal route are available at the end.")

browser_mock("Certificate & TRAQOM Survey (Mandatory)",
             "https://lms-tms.tertiaryinfotech.com",
             [("Log in", "with your registered email"), ("Course feedback", "complete the TRAQOM survey"),
              ("Digital attendance", "Assessment attendance must be submitted"),
              ("Certificate", "download once assessed Competent")],
             ["Log in to lms-tms.tertiaryinfotech.com",
              "Complete the mandatory TRAQOM survey",
              "Submit the Assessment digital attendance",
              "Download your Certificate of Achievement",
              "Your OpenCert from SSG follows by email"],
             kicker="ADMIN · TRAQOM & CERTIFICATION", accent=GREEN,
             summary="TRAQOM completion and 75% attendance are both mandatory for WSQ funding.")

tile_grid("Support", [
    ("Email", "enquiry@tertiaryinfotech.com"),
    ("Telephone", "+65 6100 0613"),
    ("Website", "www.tertiarycourses.com.sg"),
    ("LMS / TMS", "https://lms-tms.tertiaryinfotech.com"),
], kicker="WE'RE HERE AFTER THE COURSE TOO", cols=2, size=14, accent=BLUE)

s = slide(); rect(s, 0, 0, SW, SH, WHITE)
rect(s, 0, 0, SW, Inches(0.22), BLUE); rect(s, 0, Inches(7.28), SW, Inches(0.22), TEAL)
txt(s, 0, Inches(2.9), SW, Inches(1.3), [[("Thank You!", 54, INK, True)]], align=PP_ALIGN.CENTER)
txt(s, 0, Inches(4.2), SW, Inches(0.6), [[(C.TITLE, 20, GREY, False)]], align=PP_ALIGN.CENTER)
txt(s, 0, Inches(4.75), SW, Inches(0.5),
    [[(f"{C.COURSE_CODE}  ·  {C.ORG}  ·  {C.UEN}", 13, GREY, False)]], align=PP_ALIGN.CENTER)
PAGE["n"] += 1

# ---------------------------------------------------------------- motion pass
for i, sl in enumerate(prs.slides):
    shapes = list(sl.shapes)
    is_divider = len(shapes) <= 9
    _transition(sl, "push" if is_divider else "fade", "med" if is_divider else "fast")

# ---------------------------------------------------------------- save + verify
OUTDIR = os.path.join(REPO, "courseware")
os.makedirs(OUTDIR, exist_ok=True)
OUT = os.path.join(OUTDIR, f"WSQ - Master Trainer Slides - {C.SHORT_TITLE} - {C.VERSION}.pptx")
prs.save(OUT)

with open(os.path.join(HERE, "slide_map.json"), "w") as f:
    json.dump(SLIDE_MAP, f, indent=2)

# assert every generated chart asset landed on a slide
all_charts = {f for f in os.listdir(CHARTS) if f.endswith(".png")}
unused = sorted(all_charts - USED_ASSETS)
print(f"Saved: {OUT}")
print(f"Slides: {len(prs.slides.__iter__.__self__._sldIdLst)}")
if unused:
    print("WARNING — assets generated but never placed on a slide:")
    for u in unused: print("   ", u)
else:
    print(f"All {len(all_charts)} chart assets are placed on slides.")
